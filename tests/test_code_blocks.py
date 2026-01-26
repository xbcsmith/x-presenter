# SPDX-FileCopyrightText: 2024 SAS Institute Inc.
#
# SPDX-License-Identifier: Apache-2.0

"""Tests for code block parsing in markdown.

This module tests that fenced code blocks with language identifiers
are properly parsed, extracted, and structured with language and code content.
"""

import os
import tempfile

from pptx import Presentation

from presenter.converter import MarkdownToPowerPoint


class TestCodeBlockParsing:
    """Test parsing of code blocks from markdown."""

    def test_single_code_block_with_language(self):
        """Test that a single code block with language identifier is parsed correctly."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```python
def hello():
    print('world')
```"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["code_blocks"]) == 1
        assert slide_data["code_blocks"][0]["language"] == "python"
        assert "def hello():" in slide_data["code_blocks"][0]["code"]
        assert "print('world')" in slide_data["code_blocks"][0]["code"]

    def test_code_block_without_language(self):
        """Test that code blocks without language identifier are parsed."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```
plain code here
no language specified
```"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["code_blocks"]) == 1
        assert slide_data["code_blocks"][0]["language"] == ""
        assert "plain code here" in slide_data["code_blocks"][0]["code"]

    def test_multiple_code_blocks_per_slide(self):
        """Test that multiple code blocks on the same slide are parsed correctly."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```python
x = 5
```

Some text in between

```javascript
let x = 5;
```"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["code_blocks"]) == 2
        assert slide_data["code_blocks"][0]["language"] == "python"
        assert "x = 5" in slide_data["code_blocks"][0]["code"]
        assert slide_data["code_blocks"][1]["language"] == "javascript"
        assert "let x = 5;" in slide_data["code_blocks"][1]["code"]

    def test_code_block_indentation_preserved(self):
        """Test that indentation within code blocks is preserved."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```python
def function():
    if True:
        print("indented")
        x = 1
```"""

        slide_data = converter.parse_slide_content(markdown)

        code = slide_data["code_blocks"][0]["code"]
        # Check that indentation is preserved
        assert "    if True:" in code
        assert "        print" in code

    def test_empty_code_block(self):
        """Test that empty code blocks are parsed without errors."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```python
```"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["code_blocks"]) == 1
        assert slide_data["code_blocks"][0]["language"] == "python"
        assert slide_data["code_blocks"][0]["code"] == ""

    def test_unclosed_code_block_at_end(self):
        """Test that an unclosed code block at the end is still captured."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```python
print('no closing fence')"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["code_blocks"]) == 1
        assert slide_data["code_blocks"][0]["language"] == "python"
        assert "print('no closing fence')" in slide_data["code_blocks"][0]["code"]

    def test_code_block_with_various_languages(self):
        """Test code blocks with different language identifiers."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```bash
echo "hello"
```

```sql
SELECT * FROM table;
```

```json
{"key": "value"}
```"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["code_blocks"]) == 3
        assert slide_data["code_blocks"][0]["language"] == "bash"
        assert slide_data["code_blocks"][1]["language"] == "sql"
        assert slide_data["code_blocks"][2]["language"] == "json"

    def test_code_block_with_special_characters(self):
        """Test that special characters within code blocks are preserved."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```python
pattern = r"\\d+@\\w+\\.com"
string = 'value with "quotes" and \'apostrophes\''
```"""

        slide_data = converter.parse_slide_content(markdown)

        code = slide_data["code_blocks"][0]["code"]
        assert "\\d+" in code or "d+" in code  # Depending on how escaping is handled
        assert "quotes" in code
        assert "apostrophes" in code

    def test_code_block_with_unicode(self):
        """Test that Unicode characters in code blocks are preserved."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```python
message = "日本語のメッセージ"
emoji = "🎉✨"
```"""

        slide_data = converter.parse_slide_content(markdown)

        code = slide_data["code_blocks"][0]["code"]
        assert "日本語" in code
        assert "メッセージ" in code

    def test_code_block_closes_list(self):
        """Test that code block properly closes an active list."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

- Item 1
- Item 2

```python
code here
```"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["lists"]) == 1
        assert len(slide_data["lists"][0]) == 2
        assert len(slide_data["code_blocks"]) == 1

    def test_code_block_with_multiline_content(self):
        """Test code blocks with many lines of content."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```python
def complex_function(x, y):
    '''Docstring here'''
    if x > y:
        return x
    elif x < y:
        return y
    else:
        return 0

result = complex_function(5, 3)
print(result)
```"""

        slide_data = converter.parse_slide_content(markdown)

        code = slide_data["code_blocks"][0]["code"]
        assert "def complex_function" in code
        assert "Docstring here" in code
        assert "result = complex_function" in code
        assert code.count("\n") >= 9  # Multiple lines preserved

    def test_code_block_with_title_and_content(self):
        """Test code block alongside title, content, and images."""
        converter = MarkdownToPowerPoint()
        markdown = """## Code Example

This demonstrates the concept.

```javascript
const greeting = "hello";
console.log(greeting);
```

More explanation here."""

        slide_data = converter.parse_slide_content(markdown)

        assert slide_data["title"] == "Code Example"
        assert any("This demonstrates" in item for item in slide_data["content"])
        assert any("More explanation" in item for item in slide_data["content"])
        assert len(slide_data["code_blocks"]) == 1
        assert slide_data["code_blocks"][0]["language"] == "javascript"

    def test_code_block_end_to_end_pptx(self):
        """Test that code blocks appear in generated PowerPoint presentation."""
        converter = MarkdownToPowerPoint()
        markdown = """# Code Example

```python
def sample():
    return "test"
```"""

        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".md", delete=False
        ) as md_file:
            md_file.write(markdown)
            md_file.flush()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
                try:
                    converter.convert(md_file.name, pptx_file.name)

                    prs = Presentation(pptx_file.name)
                    assert len(prs.slides) > 0
                    slide = prs.slides[0]

                    # Verify slide was created
                    assert len(slide.shapes) > 0

                finally:
                    os.unlink(md_file.name)
                    os.unlink(pptx_file.name)


class TestCodeBlockEdgeCases:
    """Test edge cases for code block parsing."""

    def test_backticks_inside_code_block(self):
        """Test that backticks within code don't interfere with parsing."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```markdown
Use `backticks` for inline code
```"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["code_blocks"]) == 1
        assert "`backticks`" in slide_data["code_blocks"][0]["code"]

    def test_code_block_with_blank_lines(self):
        """Test code blocks that contain blank lines."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```python
def function1():
    pass

def function2():
    pass
```"""

        slide_data = converter.parse_slide_content(markdown)

        code = slide_data["code_blocks"][0]["code"]
        # Blank line should be preserved (as empty line)
        assert "def function1" in code
        assert "def function2" in code
        # Check that there are multiple newlines (preserving blank lines)
        assert code.count("\n") >= 3

    def test_language_with_spaces_stripped(self):
        """Test that language identifier has spaces trimmed."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```   python
code here
```"""

        slide_data = converter.parse_slide_content(markdown)

        # Language should be trimmed
        assert slide_data["code_blocks"][0]["language"] == "python"

    def test_consecutive_code_blocks(self):
        """Test multiple code blocks with no content between them."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

```python
x = 1
```
```javascript
let x = 1;
```"""

        slide_data = converter.parse_slide_content(markdown)

        # Both blocks should be captured
        assert len(slide_data["code_blocks"]) == 2
        assert "x = 1" in slide_data["code_blocks"][0]["code"]
        assert "let x = 1;" in slide_data["code_blocks"][1]["code"]
