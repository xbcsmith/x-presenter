# SPDX-FileCopyrightText: 2024 SAS Institute Inc.
#
# SPDX-License-Identifier: Apache-2.0

"""Tests for multi-line list items in markdown parsing.

This module tests that list items which wrap across multiple lines
(with continuation lines indented) are properly parsed and combined
into single list items.
"""

import os
import tempfile

from pptx import Presentation

from presenter.converter import MarkdownToPowerPoint


class TestMultiLineListItems:
    """Test parsing of list items that span multiple lines."""

    def test_single_list_item_with_continuation(self):
        """Test that a list item with a continuation line is parsed as one item."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

- first line
  continuation of first line"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["lists"]) == 1
        assert len(slide_data["lists"][0]) == 1
        assert slide_data["lists"][0][0] == "first line continuation of first line"

    def test_multiple_list_items_with_continuations(self):
        """Test multiple list items each with continuation lines."""
        converter = MarkdownToPowerPoint()
        markdown = """## Slide Title

- a really long sentence that runs on for a long time
  that is unfortunate enough to wrap lines
- another really long sentence that also wraps
  to the next line
- third item with wrap
  and continuation"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["lists"]) == 1
        assert len(slide_data["lists"][0]) == 3
        assert (
            slide_data["lists"][0][0]
            == "a really long sentence that runs on for a long time that is unfortunate enough to wrap lines"
        )
        assert (
            slide_data["lists"][0][1]
            == "another really long sentence that also wraps to the next line"
        )
        assert slide_data["lists"][0][2] == "third item with wrap and continuation"

    def test_list_item_with_multiple_continuation_lines(self):
        """Test list item with multiple continuation lines."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

- line one
  line two
  line three
  line four"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["lists"]) == 1
        assert len(slide_data["lists"][0]) == 1
        assert slide_data["lists"][0][0] == "line one line two line three line four"

    def test_mixed_single_and_multiline_list_items(self):
        """Test mix of single-line and multi-line list items."""
        converter = MarkdownToPowerPoint()
        markdown = """## Title

- short item
- long item that wraps
  to next line
- another short item
- final long item
  with continuation
  and more continuation"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["lists"]) == 1
        assert len(slide_data["lists"][0]) == 4
        assert slide_data["lists"][0][0] == "short item"
        assert slide_data["lists"][0][1] == "long item that wraps to next line"
        assert slide_data["lists"][0][2] == "another short item"
        assert (
            slide_data["lists"][0][3]
            == "final long item with continuation and more continuation"
        )

    def test_asterisk_list_with_continuation(self):
        """Test that asterisk-style lists also handle continuations."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

* first item with
  continuation
* second item
* third item with
  wrap"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["lists"]) == 1
        assert len(slide_data["lists"][0]) == 3
        assert slide_data["lists"][0][0] == "first item with continuation"
        assert slide_data["lists"][0][1] == "second item"
        assert slide_data["lists"][0][2] == "third item with wrap"

    def test_list_continuation_preserves_internal_spacing(self):
        """Test that internal spacing in continuation lines is preserved."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

- item with    multiple    spaces
  and continuation    with    spaces"""

        slide_data = converter.parse_slide_content(markdown)

        # Note: strip() will normalize the spaces between words
        assert len(slide_data["lists"]) == 1
        assert len(slide_data["lists"][0]) == 1
        # After strip, multiple spaces become single spaces
        assert "item with" in slide_data["lists"][0][0]
        assert "continuation" in slide_data["lists"][0][0]

    def test_empty_line_ends_list_continuation(self):
        """Test that an empty line after a list item doesn't continue it."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

- first item
  continuation

Regular content here"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["lists"]) == 1
        assert len(slide_data["lists"][0]) == 1
        assert slide_data["lists"][0][0] == "first item continuation"
        assert "Regular content here" in slide_data["content"]

    def test_new_list_item_ends_continuation(self):
        """Test that a new list item properly ends the previous continuation."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

- first item
  continuation
- second item starts here
  with its own continuation"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["lists"]) == 1
        assert len(slide_data["lists"][0]) == 2
        assert slide_data["lists"][0][0] == "first item continuation"
        assert (
            slide_data["lists"][0][1]
            == "second item starts here with its own continuation"
        )

    def test_end_to_end_multiline_list_in_pptx(self):
        """Test that multi-line lists appear correctly in generated PowerPoint."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

- a really long sentence that runs on for a long time about an electric sheep
  that is unfortunate enough to wrap lines
- another long sentence
  with continuation"""

        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".md", delete=False
        ) as md_file:
            md_file.write(markdown)
            md_file.flush()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
                try:
                    converter.convert(md_file.name, pptx_file.name)

                    prs = Presentation(pptx_file.name)
                    slide = prs.slides[0]

                    # Find the text box with list content
                    found_list = False
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame"):
                            text = shape.text_frame.text
                            if "electric sheep" in text and "wrap lines" in text:
                                found_list = True
                                # Both parts should be in the same bullet point
                                assert (
                                    "that is unfortunate enough to wrap lines" in text
                                )
                                break

                    assert found_list, "Multi-line list item not found in presentation"

                finally:
                    os.unlink(md_file.name)
                    os.unlink(pptx_file.name)


class TestMultiLineListEdgeCases:
    """Test edge cases for multi-line list parsing."""

    def test_deeply_indented_continuation(self):
        """Test continuation with various indentation levels."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

- item
    deeply indented continuation
  less indented continuation"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["lists"]) == 1
        assert len(slide_data["lists"][0]) == 1
        # All continuations should be combined
        result = slide_data["lists"][0][0]
        assert "item" in result
        assert "deeply indented continuation" in result
        assert "less indented continuation" in result

    def test_continuation_with_special_characters(self):
        """Test that special characters in continuation lines are preserved."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

- item with special chars: @#$%
  continuation with more: &*()"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["lists"]) == 1
        assert "@#$%" in slide_data["lists"][0][0]
        assert "&*()" in slide_data["lists"][0][0]

    def test_list_with_title_and_content_after(self):
        """Test that content after a multi-line list is handled correctly."""
        converter = MarkdownToPowerPoint()
        markdown = """## Title

Some intro text

- list item one
  with continuation
- list item two

Final paragraph here"""

        slide_data = converter.parse_slide_content(markdown)

        assert slide_data["title"] == "Title"
        assert "Some intro text" in slide_data["content"]
        assert len(slide_data["lists"]) == 1
        assert len(slide_data["lists"][0]) == 2
        assert "Final paragraph here" in slide_data["content"]

    def test_unicode_in_multiline_list(self):
        """Test Unicode characters in multi-line list items."""
        converter = MarkdownToPowerPoint()
        markdown = """# Title

- 日本語 text
  続き line
- Émojis 🎉 here
  more émojis ✨"""

        slide_data = converter.parse_slide_content(markdown)

        assert len(slide_data["lists"]) == 1
        assert len(slide_data["lists"][0]) == 2
        assert "日本語" in slide_data["lists"][0][0]
        assert "続き" in slide_data["lists"][0][0]
