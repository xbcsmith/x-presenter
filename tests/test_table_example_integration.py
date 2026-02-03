#!/usr/bin/env python3
"""
Focused integration tests for the table example markdown in testdata.

These tests:
- Use the example markdown at testdata/content/table_examples.md
- Convert it to a .pptx using the public converter API
- Open the generated presentation with python-pptx and assert native table
  shapes exist and header text is preserved.

These are integration tests and will be skipped early if `python-pptx` is not
available in the test environment.
"""

import os
import sys
import tempfile
from typing import List

import pytest

# Ensure src/ is importable when running tests directly
sys.path.insert(
    0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "src")
)

# Skip tests if python-pptx isn't installed
pytest.importorskip("pptx")

from pptx import Presentation  # type: ignore

from presenter.converter import MarkdownToPowerPoint  # type: ignore


def _count_tables_in_presentation(pres: Presentation) -> int:
    """Count table shapes across all slides in a presentation."""
    count = 0
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "table"):
                try:
                    if shape.table is not None:
                        count += 1
                except Exception:
                    continue
    return count


class TestTableExampleIntegration:
    """Integration tests that verify the `table_examples.md` fixture converts to PPTX tables."""

    def test_table_example_renders_native_tables(self) -> None:
        """Convert the example markdown and assert native table shapes are present."""
        # Resolve the path to the testdata markdown
        md_path = os.path.join(
            os.path.dirname(os.path.abspath(__file__)),
            "testdata",
            "content",
            "table_examples.md",
        )
        assert os.path.exists(md_path), f"Test data not found: {md_path}"

        converter = MarkdownToPowerPoint()

        with tempfile.TemporaryDirectory() as td:
            pptx_path = os.path.join(td, "table_examples.pptx")

            # Perform conversion
            converter.convert(md_path, pptx_path)

            # Load the generated presentation
            pres = Presentation(pptx_path)

            # Expect at least one table in the generated presentation
            table_count = _count_tables_in_presentation(pres)
            assert table_count >= 1, f"Expected at least 1 table, found {table_count}"

    def test_table_example_preserves_headers(self) -> None:
        """Verify at least one rendered table contains expected header text (Name/Role/Location or Metric)."""
        md_path = os.path.join(
            os.path.dirname(os.path.abspath(__file__)),
            "testdata",
            "content",
            "table_examples.md",
        )
        assert os.path.exists(md_path), f"Test data not found: {md_path}"

        converter = MarkdownToPowerPoint()

        with tempfile.TemporaryDirectory() as td:
            pptx_path = os.path.join(td, "table_examples_headers.pptx")
            converter.convert(md_path, pptx_path)
            pres = Presentation(pptx_path)

            header_found = False
            expected_headers = {"Name", "Role", "Location", "Metric", "Value"}

            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "table"):
                        try:
                            tbl = shape.table
                            # Read first row texts
                            first_row_texts: List[str] = [
                                tbl.cell(0, c).text.strip()
                                for c in range(tbl.columns.__len__())
                            ]
                            if any(h in expected_headers for h in first_row_texts):
                                header_found = True
                                break
                        except Exception:
                            continue
                if header_found:
                    break

            assert header_found, (
                "Expected header text (e.g. Name/Role/Metric) in at least one rendered table"
            )
