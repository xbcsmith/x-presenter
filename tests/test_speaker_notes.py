"""
Tests for speaker notes functionality in markdown to PowerPoint conversion.

SPDX-FileCopyrightText: 2025 SAS Institute Inc.
SPDX-License-Identifier: Apache-2.0
"""

import os
import tempfile

from pptx import Presentation

from presenter.converter import MarkdownToPowerPoint


class TestSpeakerNotes:
    """Test suite for HTML comment extraction and speaker notes conversion."""

    def test_parse_single_comment_as_speaker_note(self):
        """Test parsing a single HTML comment as speaker note."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide

This is content

<!-- This is a speaker note -->

- Bullet 1
- Bullet 2
"""
        slide_data = converter.parse_slide_content(markdown)

        assert slide_data["title"] == "Test Slide"
        assert "This is content" in slide_data["content"]
        assert slide_data["speaker_notes"] == "This is a speaker note"
        assert len(slide_data["lists"]) == 1
        assert slide_data["lists"][0] == ["Bullet 1", "Bullet 2"]

    def test_parse_multiple_comments_as_speaker_notes(self):
        """Test parsing multiple HTML comments combined into speaker notes."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide

<!-- First note -->

Content here

<!-- Second note -->

More content

<!-- Third note -->
"""
        slide_data = converter.parse_slide_content(markdown)

        assert slide_data["title"] == "Test Slide"
        expected_notes = "First note\n\nSecond note\n\nThird note"
        assert slide_data["speaker_notes"] == expected_notes

    def test_parse_multiline_comment_as_speaker_note(self):
        """Test parsing multiline HTML comment as speaker note."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide

<!--
This is a multiline
speaker note
with several lines
-->

Content
"""
        slide_data = converter.parse_slide_content(markdown)

        assert slide_data["title"] == "Test Slide"
        expected_notes = "This is a multiline\nspeaker note\nwith several lines"
        assert slide_data["speaker_notes"] == expected_notes

    def test_comment_not_in_slide_content(self):
        """Test that HTML comments are removed from slide content."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide

Before comment
<!-- This should not appear in content -->
After comment
"""
        slide_data = converter.parse_slide_content(markdown)

        assert slide_data["title"] == "Test Slide"
        assert "Before comment" in slide_data["content"]
        assert "After comment" in slide_data["content"]
        assert "This should not appear" not in slide_data["content"]
        assert "<!--" not in "\n".join(slide_data["content"])
        assert "-->" not in "\n".join(slide_data["content"])

    def test_empty_comment_ignored(self):
        """Test that empty HTML comments are ignored."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide

<!--  -->
<!-- -->

Content
"""
        slide_data = converter.parse_slide_content(markdown)

        assert slide_data["speaker_notes"] == ""

    def test_comment_with_special_characters(self):
        """Test HTML comment with special characters in speaker notes."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide

<!-- Remember: Use <emphasis> and don't forget the "quotes" & symbols! -->

Content
"""
        slide_data = converter.parse_slide_content(markdown)

        expected = 'Remember: Use <emphasis> and don\'t forget the "quotes" & symbols!'
        assert slide_data["speaker_notes"] == expected

    def test_no_comments_empty_speaker_notes(self):
        """Test that slides without comments have empty speaker notes."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide

Just regular content
- Bullet 1
- Bullet 2
"""
        slide_data = converter.parse_slide_content(markdown)

        assert slide_data["speaker_notes"] == ""

    def test_speaker_notes_added_to_presentation(self):
        """Test that speaker notes are added to PowerPoint slides."""
        converter = MarkdownToPowerPoint()
        slide_data = {
            "title": "Test Slide",
            "content": ["Some content"],
            "lists": [],
            "images": [],
            "speaker_notes": "Important note for presenter",
        }

        converter.add_slide_to_presentation(slide_data)

        # Verify slide was created
        assert len(converter.presentation.slides) == 1

        # Check speaker notes
        slide = converter.presentation.slides[0]
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text

        assert notes_text == "Important note for presenter"

    def test_speaker_notes_with_no_notes_field(self):
        """Test handling slide data without speaker_notes field."""
        converter = MarkdownToPowerPoint()
        slide_data = {
            "title": "Test Slide",
            "content": ["Some content"],
            "lists": [],
            "images": [],
        }

        # Should not raise an error
        converter.add_slide_to_presentation(slide_data)

        assert len(converter.presentation.slides) == 1

    def test_speaker_notes_with_empty_string(self):
        """Test handling slide data with empty speaker notes string."""
        converter = MarkdownToPowerPoint()
        slide_data = {
            "title": "Test Slide",
            "content": ["Some content"],
            "lists": [],
            "images": [],
            "speaker_notes": "",
        }

        converter.add_slide_to_presentation(slide_data)

        assert len(converter.presentation.slides) == 1

    def test_end_to_end_speaker_notes_in_pptx(self):
        """Test full conversion from markdown with comments to PPTX with notes."""
        markdown_content = """# Slide 1

This is the first slide

<!-- Remember to introduce yourself -->

- Point 1
- Point 2

---

# Slide 2

Second slide content

<!--
Multi-line note:
- Emphasize key metrics
- Show enthusiasm
-->

- Bullet A
- Bullet B

---

# Slide 3

Final slide

No speaker notes here
"""

        with tempfile.TemporaryDirectory() as temp_dir:
            # Create markdown file
            md_file = os.path.join(temp_dir, "test.md")
            with open(md_file, "w", encoding="utf-8") as f:
                f.write(markdown_content)

            # Convert to PowerPoint
            output_file = os.path.join(temp_dir, "output.pptx")
            converter = MarkdownToPowerPoint()
            converter.convert(md_file, output_file)

            # Verify output file exists
            assert os.path.exists(output_file)

            # Load the presentation and verify notes
            prs = Presentation(output_file)
            assert len(prs.slides) == 3

            # Check slide 1 notes
            slide1_notes = prs.slides[0].notes_slide.notes_text_frame.text
            assert slide1_notes == "Remember to introduce yourself"

            # Check slide 2 notes
            slide2_notes = prs.slides[1].notes_slide.notes_text_frame.text
            expected = "Multi-line note:\n- Emphasize key metrics\n- Show enthusiasm"
            assert slide2_notes == expected

            # Check slide 3 has no notes (empty string)
            slide3_notes = prs.slides[2].notes_slide.notes_text_frame.text
            assert slide3_notes == ""

    def test_comment_between_list_items(self):
        """Test HTML comment placed between list items."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide

- Item 1
<!-- Note about item 1 -->
- Item 2
- Item 3
"""
        slide_data = converter.parse_slide_content(markdown)

        # Comment should be in notes
        assert slide_data["speaker_notes"] == "Note about item 1"

        # List should still be intact
        assert len(slide_data["lists"]) == 1
        assert slide_data["lists"][0] == ["Item 1", "Item 2", "Item 3"]

    def test_comment_after_title(self):
        """Test HTML comment immediately after title."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide
<!-- This note is right after the title -->

Content here
"""
        slide_data = converter.parse_slide_content(markdown)

        assert slide_data["title"] == "Test Slide"
        assert slide_data["speaker_notes"] == "This note is right after the title"
        assert "Content here" in slide_data["content"]

    def test_comment_with_markdown_inside(self):
        """Test HTML comment containing markdown syntax (should be preserved)."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide

<!-- Remember to emphasize **bold** points and mention `code` examples -->

Content
"""
        slide_data = converter.parse_slide_content(markdown)

        expected = "Remember to emphasize **bold** points and mention `code` examples"
        assert slide_data["speaker_notes"] == expected

    def test_nested_comments_not_supported(self):
        """Test that nested HTML comments are handled (outer comment wins)."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide

<!-- Outer comment <!-- inner --> still outer -->

Content
"""
        slide_data = converter.parse_slide_content(markdown)

        # Due to regex greedy matching, this will capture up to first -->
        # This is expected HTML comment behavior
        assert "Outer comment" in slide_data["speaker_notes"]

    def test_multiple_slides_with_different_notes(self):
        """Test multiple slides each with different speaker notes."""
        markdown_content = """# Slide A
<!-- Note A -->
Content A

---

# Slide B
<!-- Note B1 -->
Content B
<!-- Note B2 -->

---

# Slide C
Content C without notes
"""
        converter = MarkdownToPowerPoint()
        slides = converter.parse_markdown_slides(markdown_content)

        assert len(slides) == 3

        slide_a = converter.parse_slide_content(slides[0])
        assert slide_a["speaker_notes"] == "Note A"

        slide_b = converter.parse_slide_content(slides[1])
        assert slide_b["speaker_notes"] == "Note B1\n\nNote B2"

        slide_c = converter.parse_slide_content(slides[2])
        assert slide_c["speaker_notes"] == ""


class TestSpeakerNotesEdgeCases:
    """Test edge cases and error conditions for speaker notes."""

    def test_comment_only_slide(self):
        """Test slide containing only HTML comments (no visible content)."""
        converter = MarkdownToPowerPoint()
        markdown = """<!-- Just a note, no visible content -->"""

        slide_data = converter.parse_slide_content(markdown)

        assert slide_data["title"] == ""
        assert slide_data["content"] == []
        assert slide_data["speaker_notes"] == "Just a note, no visible content"

    def test_very_long_speaker_note(self):
        """Test handling of very long speaker notes."""
        converter = MarkdownToPowerPoint()
        long_note = "This is a very long note. " * 100
        markdown = f"""# Test Slide

<!-- {long_note} -->

Content
"""
        slide_data = converter.parse_slide_content(markdown)

        assert slide_data["speaker_notes"] == long_note.strip()
        assert len(slide_data["speaker_notes"]) > 1000

    def test_comment_with_newlines_preserved(self):
        """Test that newlines in comments are preserved."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide

<!--
Line 1
Line 2

Line 4 (with blank line above)
-->

Content
"""
        slide_data = converter.parse_slide_content(markdown)

        assert "Line 1" in slide_data["speaker_notes"]
        assert "Line 2" in slide_data["speaker_notes"]
        assert "Line 4" in slide_data["speaker_notes"]
        # Newlines should be preserved
        assert "\n" in slide_data["speaker_notes"]

    def test_unicode_in_speaker_notes(self):
        """Test Unicode characters in speaker notes."""
        converter = MarkdownToPowerPoint()
        markdown = """# Test Slide

<!-- 你好 world! Émoji: 🎉 Math: ∑∫ -->

Content
"""
        slide_data = converter.parse_slide_content(markdown)

        assert "你好 world!" in slide_data["speaker_notes"]
        assert "🎉" in slide_data["speaker_notes"]
        assert "∑∫" in slide_data["speaker_notes"]
