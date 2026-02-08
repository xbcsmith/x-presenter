#!/usr/bin/env python3
"""
Integration tests (end-to-end) for Markdown table support.

These tests exercise the converter end-to-end:
- Write markdown files containing tables and other content
- Run the converter to produce a .pptx file
- Open the generated presentation and assert that tables and other
  elements have been rendered into the slides

Note: These are integration tests that require `python-pptx` to be
available in the test environment. If `python-pptx` is not installed,
the tests in this module are skipped early.
"""

import os
import sys
import tempfile
import textwrap
from typing import List

import pytest

# Make sure src is importable when running tests directly
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "src"))

# Skip the tests in this module early if python-pptx is not available.
pytest.importorskip("pptx")

from pptx import Presentation  # type: ignore

from presenter.converter import MarkdownToPowerPoint  # type: ignore


def _count_tables_in_presentation(pres: Presentation) -> int:
    """Count table shapes across all slides in a presentation."""
    count = 0
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "table"):
                # Some shapes may have a .table attribute even if empty; check truthiness
                try:
                    if shape.table is not None:
                        count += 1
                except Exception:
                    # be defensive: if accessing table property raises, skip
                    continue
    return count


def _presentation_has_image(pres: Presentation) -> bool:
    """Return True if any slide in presentation contains a picture shape."""
    for slide in pres.slides:
        for shape in slide.shapes:
            # Picture shapes expose an 'image' attribute in python-pptx
            if hasattr(shape, "image"):
                return True
    return False


def _presentation_has_speaker_notes(pres: Presentation, expected_substring: str) -> bool:
    """Check if any slide's notes contain the expected substring."""
    for slide in pres.slides:
        try:
            notes = slide.notes_slide.notes_text_frame.text
            if notes and expected_substring in notes:
                return True
        except Exception:
            continue
    return False


class TestTableIntegrationEndToEnd:
    """Integration tests that verify tables are rendered in PPTX output."""

    def test_full_presentation_with_tables(self) -> None:
        """
        Create a small multi-slide markdown document with multiple tables,
        convert to PPTX, and assert that tables exist and headers preserved.
        """
        converter = MarkdownToPowerPoint()

        markdown = textwrap.dedent(
            """\
            # Project Overview

            | Metric | Value |
            |-------:|:-----:|
            | Users  | 1024  |
            | Uptime | 99.9% |

            ---

            ## Team Directory

            Team list and table below:

            | Name  | Role      | Location   |
            |:------|:---------:|----------:|
            | Alice | Engineer  | London     |
            | Bob   | Designer  | New York   |

            """
        )

        # Write markdown to a temp file and convert
        with tempfile.TemporaryDirectory() as td:
            md_path = os.path.join(td, "slides.md")
            pptx_path = os.path.join(td, "slides.pptx")
            with open(md_path, "w", encoding="utf-8") as f:
                f.write(markdown)

            # Perform conversion
            converter.convert(md_path, pptx_path)

            # Load generated presentation
            pres = Presentation(pptx_path)

            # There should be two slides (title + content slide)
            assert len(pres.slides) >= 2

            # Count tables - expect at least two tables (one per slide)
            table_count = _count_tables_in_presentation(pres)
            assert table_count >= 2, f"Expected >=2 tables, found {table_count}"

            # Verify header presence for one of the tables (search across all tables)
            header_found = False
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "table"):
                        try:
                            tbl = shape.table
                            # check first row cells contain any of the known header labels
                            first_row_texts: List[str] = [
                                tbl.cell(0, c).text.strip() for c in range(tbl.columns.__len__())
                            ]
                            if any(h in ("Metric", "Name", "Role") for h in first_row_texts):
                                header_found = True
                                break
                        except Exception:
                            continue
                if header_found:
                    break

            assert header_found, "Expected header text in at least one rendered table"

    def test_table_with_all_content_types(self) -> None:
        """
        Create a markdown slide that mixes a table with lists, inline code,
        an image, and speaker notes. Convert and assert that all these elements
        are present in the generated presentation.
        """
        converter = MarkdownToPowerPoint()

        # Create a tiny 1x1 PNG image (binary) to reference from markdown.
        # This avoids a dependency on Pillow for tests; the bytes below are a valid PNG.
        tiny_png_bytes = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89"
            b"\x00\x00\x00\nIDATx\xdac\xf8\x0f\x00\x01\x01\x01\x00"
            b"\x18\xdd\x03\x16\x00\x00\x00\x00IEND\xaeB`\x82"
        )

        markdown = textwrap.dedent(
            """\
            ## Mixed Content Slide

            This slide mixes inline `code`, a bullet list, an image, and a table.

            - First bullet
            - Second bullet with `inline_code`

            | Key | Description |
            |:----|------------:|
            | CPU | 4 cores     |
            | RAM | 16GB        |

            ![logo](./logo.png)

            <!-- NOTE: This slide contains sensitive talking points -->
            """
        )

        with tempfile.TemporaryDirectory() as td:
            md_path = os.path.join(td, "mixed.md")
            img_path = os.path.join(td, "logo.png")
            pptx_path = os.path.join(td, "mixed.pptx")

            # Write image bytes
            with open(img_path, "wb") as imgf:
                imgf.write(tiny_png_bytes)

            # Write markdown
            with open(md_path, "w", encoding="utf-8") as f:
                f.write(markdown)

            # Convert
            converter.convert(md_path, pptx_path)

            # Load presentation and assert elements
            pres = Presentation(pptx_path)

            # Expect at least one table
            assert _count_tables_in_presentation(pres) >= 1

            # Expect at least one image (the tiny png we inserted)
            assert _presentation_has_image(pres), "Expected an image shape in presentation"

            # Expect speaker notes to contain the NOTE text
            assert _presentation_has_speaker_notes(pres, "sensitive talking points"), (
                "Expected speaker notes to contain 'sensitive talking points'"
            )
