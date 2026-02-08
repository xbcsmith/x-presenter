# SPDX-FileCopyrightText: 2024 SAS Institute Inc.
#
# SPDX-License-Identifier: Apache-2.0

"""Tests for title slide vs content slide layout functionality.

This module tests that the first slide with a single # header uses the
Title Slide layout (centered title) while subsequent slides use the
Title and Content layout.
"""

import os
import tempfile

from pptx import Presentation

from presenter.converter import MarkdownToPowerPoint


class TestTitleSlideLayout:
    """Test that first slide uses Title Slide layout when starting with single #."""

    def test_first_slide_with_single_hash_is_title_slide(self):
        """Test that first slide starting with # uses title slide layout."""
        converter = MarkdownToPowerPoint()
        markdown = "# Title Slide\n---\n## Content Slide"
        slides = converter.parse_markdown_slides(markdown)

        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as md_file:
            md_file.write(markdown)
            md_file.flush()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
                try:
                    converter.convert(md_file.name, pptx_file.name)

                    # Load presentation and check layouts
                    prs = Presentation(pptx_file.name)
                    assert len(prs.slides) == 2

                    # First slide should use layout 0 (Title Slide)
                    first_slide = prs.slides[0]
                    assert first_slide.slide_layout.name == "Title Slide"

                    # Second slide should use layout 1 (Title and Content)
                    second_slide = prs.slides[1]
                    assert second_slide.slide_layout.name == "Title and Content"

                finally:
                    os.unlink(md_file.name)
                    os.unlink(pptx_file.name)

    def test_first_slide_with_double_hash_uses_content_layout(self):
        """Test that first slide starting with ## uses content layout, not title layout."""
        converter = MarkdownToPowerPoint()
        markdown = "## Subtitle Slide\n---\n## Another Slide"
        slides = converter.parse_markdown_slides(markdown)

        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as md_file:
            md_file.write(markdown)
            md_file.flush()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
                try:
                    converter.convert(md_file.name, pptx_file.name)

                    # Load presentation and check layouts
                    prs = Presentation(pptx_file.name)
                    assert len(prs.slides) == 2

                    # Both slides should use layout 1 (Title and Content)
                    first_slide = prs.slides[0]
                    assert first_slide.slide_layout.name == "Title and Content"

                    second_slide = prs.slides[1]
                    assert second_slide.slide_layout.name == "Title and Content"

                finally:
                    os.unlink(md_file.name)
                    os.unlink(pptx_file.name)

    def test_first_slide_without_title_uses_content_layout(self):
        """Test that first slide without title uses content layout."""
        converter = MarkdownToPowerPoint()
        markdown = "Regular content\n---\n# Second Slide"
        slides = converter.parse_markdown_slides(markdown)

        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as md_file:
            md_file.write(markdown)
            md_file.flush()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
                try:
                    converter.convert(md_file.name, pptx_file.name)

                    # Load presentation and check layouts
                    prs = Presentation(pptx_file.name)
                    assert len(prs.slides) == 2

                    # Both slides should use layout 1 (Title and Content)
                    first_slide = prs.slides[0]
                    assert first_slide.slide_layout.name == "Title and Content"

                    second_slide = prs.slides[1]
                    assert second_slide.slide_layout.name == "Title and Content"

                finally:
                    os.unlink(md_file.name)
                    os.unlink(pptx_file.name)

    def test_add_slide_with_is_title_slide_true(self):
        """Test add_slide_to_presentation with is_title_slide=True."""
        converter = MarkdownToPowerPoint()
        slide_data = {
            "title": "Main Title",
            "content": [],
            "lists": [],
            "images": [],
            "speaker_notes": "",
        }

        converter.add_slide_to_presentation(slide_data, is_title_slide=True)

        assert len(converter.presentation.slides) == 1
        slide = converter.presentation.slides[0]
        assert slide.slide_layout.name == "Title Slide"

    def test_add_slide_with_is_title_slide_false(self):
        """Test add_slide_to_presentation with is_title_slide=False."""
        converter = MarkdownToPowerPoint()
        slide_data = {
            "title": "Content Title",
            "content": ["Some content"],
            "lists": [],
            "images": [],
            "speaker_notes": "",
        }

        converter.add_slide_to_presentation(slide_data, is_title_slide=False)

        assert len(converter.presentation.slides) == 1
        slide = converter.presentation.slides[0]
        assert slide.slide_layout.name == "Title and Content"

    def test_add_slide_default_is_content_layout(self):
        """Test that is_title_slide defaults to False (content layout)."""
        converter = MarkdownToPowerPoint()
        slide_data = {
            "title": "Default Layout",
            "content": [],
            "lists": [],
            "images": [],
            "speaker_notes": "",
        }

        # Call without is_title_slide parameter
        converter.add_slide_to_presentation(slide_data)

        assert len(converter.presentation.slides) == 1
        slide = converter.presentation.slides[0]
        assert slide.slide_layout.name == "Title and Content"


class TestTitleSlideContent:
    """Test that title slide content is properly handled."""

    def test_title_slide_has_title_text(self):
        """Test that title slide receives the title text."""
        converter = MarkdownToPowerPoint()
        markdown = "# My Presentation Title"

        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as md_file:
            md_file.write(markdown)
            md_file.flush()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
                try:
                    converter.convert(md_file.name, pptx_file.name)

                    prs = Presentation(pptx_file.name)
                    slide = prs.slides[0]

                    # Check that title placeholder exists and has text
                    assert slide.shapes.title is not None
                    assert slide.shapes.title.text == "My Presentation Title"

                finally:
                    os.unlink(md_file.name)
                    os.unlink(pptx_file.name)

    def test_content_slide_has_title_text(self):
        """Test that content slide receives the title text."""
        converter = MarkdownToPowerPoint()
        markdown = "# Title\n---\n## Content Slide Title"

        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as md_file:
            md_file.write(markdown)
            md_file.flush()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
                try:
                    converter.convert(md_file.name, pptx_file.name)

                    prs = Presentation(pptx_file.name)
                    second_slide = prs.slides[1]

                    # Check that title placeholder exists and has text
                    assert second_slide.shapes.title is not None
                    assert second_slide.shapes.title.text == "Content Slide Title"

                finally:
                    os.unlink(md_file.name)
                    os.unlink(pptx_file.name)


class TestMultipleSlidesLayout:
    """Test layout assignment with multiple slides."""

    def test_multiple_slides_only_first_is_title_slide(self):
        """Test that only the first slide with # is treated as title slide."""
        converter = MarkdownToPowerPoint()
        markdown = "# Title\n---\n# Second\n---\n# Third"

        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as md_file:
            md_file.write(markdown)
            md_file.flush()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
                try:
                    converter.convert(md_file.name, pptx_file.name)

                    prs = Presentation(pptx_file.name)
                    assert len(prs.slides) == 3

                    # First slide: Title Slide layout
                    assert prs.slides[0].slide_layout.name == "Title Slide"

                    # Second and third: Title and Content layout
                    assert prs.slides[1].slide_layout.name == "Title and Content"
                    assert prs.slides[2].slide_layout.name == "Title and Content"

                finally:
                    os.unlink(md_file.name)
                    os.unlink(pptx_file.name)

    def test_all_content_slides_no_title_slide(self):
        """Test presentation with no title slide (all content slides)."""
        converter = MarkdownToPowerPoint()
        markdown = "## Slide 1\n---\n## Slide 2\n---\n## Slide 3"

        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as md_file:
            md_file.write(markdown)
            md_file.flush()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
                try:
                    converter.convert(md_file.name, pptx_file.name)

                    prs = Presentation(pptx_file.name)
                    assert len(prs.slides) == 3

                    # All slides should use Title and Content layout
                    for slide in prs.slides:
                        assert slide.slide_layout.name == "Title and Content"

                finally:
                    os.unlink(md_file.name)
                    os.unlink(pptx_file.name)


class TestTitleSlideWithSpeakerNotes:
    """Test that speaker notes work correctly with title slides."""

    def test_title_slide_with_speaker_notes(self):
        """Test that title slide can have speaker notes."""
        converter = MarkdownToPowerPoint()
        markdown = "# Title Slide\n<!-- Speaker notes for title slide -->"

        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as md_file:
            md_file.write(markdown)
            md_file.flush()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
                try:
                    converter.convert(md_file.name, pptx_file.name)

                    prs = Presentation(pptx_file.name)
                    slide = prs.slides[0]

                    # Verify it's a title slide
                    assert slide.slide_layout.name == "Title Slide"

                    # Verify speaker notes are present
                    notes_slide = slide.notes_slide
                    notes_text = notes_slide.notes_text_frame.text
                    assert "Speaker notes for title slide" in notes_text

                finally:
                    os.unlink(md_file.name)
                    os.unlink(pptx_file.name)


class TestTitleSlideWithBackground:
    """Test that title slides work correctly with background images."""

    def test_title_slide_with_background_image(self, tmp_path):
        """Test that title slide can have background image."""
        # Create a simple 1x1 pixel image for testing
        from PIL import Image

        bg_image_path = tmp_path / "background.png"
        img = Image.new("RGB", (1, 1), color="blue")
        img.save(bg_image_path)

        converter = MarkdownToPowerPoint()
        markdown = "# Title Slide"

        md_file_path = tmp_path / "test.md"
        md_file_path.write_text(markdown)

        pptx_file_path = tmp_path / "output.pptx"

        converter.convert(str(md_file_path), str(pptx_file_path), str(bg_image_path))

        prs = Presentation(str(pptx_file_path))
        slide = prs.slides[0]

        # Verify it's a title slide
        assert slide.slide_layout.name == "Title Slide"

        # Verify slide was created (background would be first shape if added)
        assert len(slide.shapes) > 0


class TestEdgeCases:
    """Test edge cases for title slide detection."""

    def test_first_slide_with_whitespace_before_hash(self):
        """Test that whitespace before # still triggers title slide layout."""
        converter = MarkdownToPowerPoint()
        # After strip() in parse, this becomes "# Title"
        markdown = "   # Title   \n---\n## Second"

        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as md_file:
            md_file.write(markdown)
            md_file.flush()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
                try:
                    converter.convert(md_file.name, pptx_file.name)

                    prs = Presentation(pptx_file.name)
                    # First slide should be title slide
                    assert prs.slides[0].slide_layout.name == "Title Slide"

                finally:
                    os.unlink(md_file.name)
                    os.unlink(pptx_file.name)

    def test_single_slide_with_hash_is_title_slide(self):
        """Test that a presentation with a single # slide uses title layout."""
        converter = MarkdownToPowerPoint()
        markdown = "# Only Slide"

        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False) as md_file:
            md_file.write(markdown)
            md_file.flush()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
                try:
                    converter.convert(md_file.name, pptx_file.name)

                    prs = Presentation(pptx_file.name)
                    assert len(prs.slides) == 1
                    assert prs.slides[0].slide_layout.name == "Title Slide"

                finally:
                    os.unlink(md_file.name)
                    os.unlink(pptx_file.name)
