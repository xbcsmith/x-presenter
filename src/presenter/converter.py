"""
Core module for converting Markdown presentations to PowerPoint.
"""

import logging
import os
from typing import TYPE_CHECKING, Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from .config import Config
from .parsers.code import (
    CODE_BLOCK_MAX_HEIGHT,
    CODE_BLOCK_MIN_HEIGHT,
    calculate_code_block_height,
    tokenize_code,
)
from .parsers.slides import parse_markdown_slides, parse_slide_content
from .parsers.tables import (
    TABLE_CELL_FONT_SIZE,
    TABLE_HEADER_BG,
    TABLE_HEADER_FONT_SIZE,
    TableParseError,
    calculate_table_dimensions,
)
from .parsers.text import parse_markdown_formatting
from .utils.colors import parse_color
from .utils.ppt_cleanup import remove_unused_placeholders

if TYPE_CHECKING:
    pass

logger = logging.getLogger(__name__)


class MarkdownToPowerPoint:
    """Convert Markdown presentations to PowerPoint format."""

    def __init__(
        self,
        background_image: Optional[str] = None,
        background_color: Optional[str] = None,
        font_color: Optional[str] = None,
        title_bg_color: Optional[str] = None,
        title_font_color: Optional[str] = None,
        code_background_color: Optional[str] = None,
    ):
        """Initialize the converter.

        Args:
            background_image: Path to background image file (optional)
            background_color: Background color for content slides (hex: RRGGBB or #RRGGBB)
            font_color: Font color for content slides (hex: RRGGBB or #RRGGBB)
            title_bg_color: Background color for title slide (hex: RRGGBB or #RRGGBB)
            title_font_color: Font color for title slide (hex: RRGGBB or #RRGGBB)
            code_background_color: Background color for code blocks (hex: RRGGBB or #RRGGBB)
        """
        self.presentation = Presentation()
        self.slide_separator = "---"
        self.background_image = background_image
        self.background_color = self._parse_color(background_color)
        self.font_color = self._parse_color(font_color)
        self.title_bg_color = self._parse_color(title_bg_color)
        self.title_font_color = self._parse_color(title_font_color)
        self.code_background_color = self._parse_color(code_background_color)
        if self.code_background_color is None:
            # Default dark background for code blocks with better contrast
            self.code_background_color = RGBColor(30, 30, 30)

        # Re-export parsing methods as instance methods for compatibility
        self.parse_markdown_slides = self._parse_markdown_slides_wrapper
        self.parse_slide_content = parse_slide_content

        # Re-export table parsing error for compatibility
        self.TableParseError = TableParseError

    def _parse_color(self, color_str: Optional[str]) -> Optional[RGBColor]:
        """Wrapper for color parsing utility."""
        return parse_color(color_str)

    def _remove_unused_placeholders(self, slide, has_title: bool, has_body_content: bool) -> None:
        """Wrapper for placeholder cleanup utility."""
        remove_unused_placeholders(slide, has_title, has_body_content)

    def _parse_markdown_slides_wrapper(self, markdown_content: str) -> List[str]:
        """Wrapper for slide parsing utility using instance separator."""
        return parse_markdown_slides(markdown_content, self.slide_separator)

    def _parse_markdown_formatting(self, text: str) -> List[Dict[str, Any]]:
        """Wrapper for text formatting utility."""
        return parse_markdown_formatting(text)

    def _render_table(self, slide, top_position, table_struct: Dict[str, Any]) -> float:
        """Render a table on the provided slide at the given top_position.

        Returns the rendered height in inches (float) so callers can advance layout.
        """
        dims = calculate_table_dimensions(table_struct)
        rows = dims["rows"]
        cols = dims["cols"]
        total_width = dims["total_width"]
        total_height = dims["total_height"]

        left = Inches(0.5)
        top = top_position
        try:
            tbl_shape = slide.shapes.add_table(rows, cols, left, top, Inches(total_width), Inches(total_height))
            tbl = tbl_shape.table
        except Exception:
            # Could not create a table; fallback to textual representation height
            fallback_height = max(rows * 0.25, 0.5)
            return fallback_height

        # Set column widths evenly
        col_width = total_width / cols
        for c in range(cols):
            tbl.columns[c].width = Inches(col_width)

        # Populate header if present
        header_offset = 0
        if table_struct.get("has_header") and table_struct.get("headers"):
            headers = table_struct["headers"]
            for c in range(cols):
                cell = tbl.cell(0, c)
                cell.text_frame.clear()
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = headers[c] if c < len(headers) else ""
                run.font.name = "Courier New"
                run.font.size = Pt(TABLE_HEADER_FONT_SIZE)
                # Apply header bg color if possible
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TABLE_HEADER_BG
                except Exception:
                    pass
            header_offset = 1

        # Populate data rows
        for r_idx, row in enumerate(table_struct.get("rows", [])):
            for c in range(cols):
                cell = tbl.cell(r_idx + header_offset, c)
                cell.text_frame.clear()
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = row[c] if c < len(row) else ""
                run.font.name = "Courier New"
                run.font.size = Pt(TABLE_CELL_FONT_SIZE)

        return total_height

    def _apply_text_formatting(self, text_frame, text: str, font_size: int = 18, color: RGBColor = None):
        """Apply markdown formatting to text frame."""
        text_frame.clear()
        p = text_frame.paragraphs[0]
        segments = parse_markdown_formatting(text)

        for segment in segments:
            run = p.add_run()
            run.text = segment["text"]
            run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color
            if segment["bold"]:
                run.font.bold = True
            if segment["italic"]:
                run.font.italic = True
            if segment["code"]:
                run.font.name = "Courier New"

    def _render_text_content(self, slide, text: str, content_type: str, top_position: Any) -> Any:
        """Render a text content block on the slide."""
        content_box = slide.shapes.add_textbox(Inches(0.5), top_position, Inches(9), Inches(0.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        p = content_frame.paragraphs[0]

        # Set spacing based on type
        if content_type.startswith("h"):
            p.space_before = Pt(6)
            p.space_after = Pt(3)
        else:
            p.space_before = Pt(3)
            p.space_after = Pt(3)

        # Apply formatting based on content type
        segments = self._parse_markdown_formatting(text)
        for segment in segments:
            run = p.add_run()
            run.text = segment["text"]
            if segment["bold"]:
                run.font.bold = True
            if segment["italic"]:
                run.font.italic = True
            if segment["code"]:
                run.font.name = "Courier New"

            # Set font size based on content type
            if content_type == "h3":
                run.font.size = Pt(22)
                run.font.bold = True
            elif content_type == "h4":
                run.font.size = Pt(20)
                run.font.bold = True
            elif content_type in ["h5", "h6"]:
                run.font.size = Pt(18)
                run.font.bold = True
            else:
                run.font.size = Pt(16)

            if self.font_color:
                run.font.color.rgb = self.font_color

        # Better height estimation to prevent overlaps
        # Assuming 9 inches width
        if content_type == "h3":
            chars_per_line = 50  # 22pt
            line_height = 0.5
        elif content_type == "h4":
            chars_per_line = 60  # 20pt
            line_height = 0.45
        elif content_type in ["h5", "h6"]:
            chars_per_line = 70  # 18pt
            line_height = 0.4
        else:
            chars_per_line = 85  # 16pt
            line_height = 0.35

        # Calculate estimated number of lines
        lines = 1
        if len(text) > chars_per_line:
            lines = (len(text) // chars_per_line) + 1

        estimated_height = lines * line_height

        return Inches(top_position.inches + estimated_height + 0.1)

    def _render_list_block(self, slide, items: List[str], top_position: Any) -> Any:
        """Render a list block on the slide."""
        # Calculate height dynamically based on content length
        total_list_height = 0
        chars_per_line = 80  # List items are usually 14-16pt
        line_height_per_item = 0.35

        for item in items:
            lines = 1
            if len(item) > chars_per_line:
                lines = (len(item) // chars_per_line) + 1
            total_list_height += lines * line_height_per_item

        list_height = max(total_list_height, 0.5)

        list_box = slide.shapes.add_textbox(Inches(1), top_position, Inches(8), Inches(list_height))
        list_frame = list_box.text_frame
        list_frame.clear()
        list_frame.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = list_frame.paragraphs[0]
            else:
                p = list_frame.add_paragraph()

            # Reduce spacing between list items
            p.space_before = Pt(0)
            p.space_after = Pt(3)

            # Add bullet point and then formatted text
            bullet_run = p.add_run()
            bullet_run.text = "• "
            bullet_run.font.size = Pt(16)
            if self.font_color:
                bullet_run.font.color.rgb = self.font_color

            # Parse and apply markdown formatting to list item
            segments = self._parse_markdown_formatting(item)
            for segment in segments:
                run = p.add_run()
                run.text = segment["text"]
                if segment["bold"]:
                    run.font.bold = True
                if segment["italic"]:
                    run.font.italic = True
                if segment["code"]:
                    run.font.name = "Courier New"
                run.font.size = Pt(14)
                if self.font_color:
                    run.font.color.rgb = self.font_color

        return Inches(top_position.inches + list_height + 0.15)

    def _render_code_block(self, slide, code_block: Dict[str, str], top_position: Any) -> Any:
        """Render a code block on the slide."""
        code_text = code_block["code"]
        language = code_block["language"]

        # Calculate height (line-based) and add small padding so text doesn't touch edges.
        # calculate_code_block_height already enforces min/max; add padding and re-apply bounds.
        block_height = calculate_code_block_height(code_text)
        padding = 0.2  # inches total (0.1 top + 0.1 bottom)
        # Ensure padding doesn't push us past the configured limits
        block_height = min(
            max(block_height + padding, CODE_BLOCK_MIN_HEIGHT),
            CODE_BLOCK_MAX_HEIGHT,
        )

        # Create textbox for code
        code_box = slide.shapes.add_textbox(
            Inches(0.5),  # Left margin
            top_position,
            Inches(9),  # Width
            Inches(block_height),
        )

        # Configure text frame
        code_frame = code_box.text_frame
        # For code blocks we want fixed shape height (calculated above) and preserved line breaks.
        # Disable auto-resizing so the shape height remains the expected block_height.
        # Also disable word wrap to preserve indentation and long lines; rely on horizontal scrolling
        # or clipping rather than undesired wrapping which can change layout unpredictably.
        code_frame.word_wrap = False
        code_frame.auto_size = MSO_AUTO_SIZE.NONE
        # Anchor text to top so the code appears at the top of the box
        code_frame.vertical_anchor = MSO_ANCHOR.TOP
        # Keep small margins so code doesn't touch the edges
        code_frame.margin_left = Inches(0.1)
        code_frame.margin_right = Inches(0.1)
        code_frame.margin_top = Inches(0.1)
        code_frame.margin_bottom = Inches(0.1)

        # Set background color (light gray or configured color)
        fill = code_box.fill
        fill.solid()
        fill.fore_color.rgb = self.code_background_color

        # Add code with syntax highlighting
        tokens = tokenize_code(code_text, language)

        # First token goes in existing paragraph
        p = code_frame.paragraphs[0]

        for token in tokens:
            text = token["text"]
            # If token contains newline(s), split on '\n' and create new paragraphs for each newline.
            # This preserves runs' colors and ensures multi-line whitespace tokens (e.g. " \n")
            # are handled correctly rather than only matching exact "\n".
            if "\n" in text:
                parts = text.split("\n")
                for idx, part in enumerate(parts):
                    if part:
                        run = p.add_run()
                        run.text = part
                        run.font.name = "Courier New"
                        run.font.size = Pt(12)
                        if token.get("color"):
                            run.font.color.rgb = token["color"]
                    # After each newline except the last, start a new paragraph
                    if idx != len(parts) - 1:
                        p = code_frame.add_paragraph()
            else:
                # Add run to current paragraph
                run = p.add_run()
                run.text = text
                run.font.name = "Courier New"
                run.font.size = Pt(12)
                if token.get("color"):
                    run.font.color.rgb = token["color"]

        return Inches(top_position.inches + block_height + 0.15)

    def _render_image(self, slide, image_info: Dict[str, str], base_path: str, top_position: Any) -> Any:
        """Render an image on the slide."""
        image_path = image_info["path"]

        # Handle relative paths
        if not os.path.isabs(image_path):
            if image_path.startswith("./"):
                image_path = image_path[2:]
            image_path = os.path.join(base_path, image_path)

        if os.path.exists(image_path):
            try:
                # Add image to slide
                slide.shapes.add_picture(image_path, Inches(2), top_position, height=Inches(3))
                return Inches(top_position.inches + 3.5)
            except Exception as e:
                print(f"Warning: Could not add image {image_path}: {e}")
        else:
            print(f"Warning: Image not found: {image_path}")

        return top_position

    def add_slide_to_presentation(
        self,
        slide_data: Dict[str, Any],
        base_path: str = "",
        is_title_slide: bool = False,
    ) -> None:
        """Add a slide to the presentation based on parsed data.

        Creates a new slide with parsed content including title, text content,
        bullet lists, images, and speaker notes. Handles background images and
        colors if configured. Automatically positions content vertically on the slide.
        Uses Title Slide layout for title slides (centered title) and Title and
        Content layout for content slides.

        Args:
            slide_data: Parsed slide data dictionary with keys:
                'title': Slide title text (str)
                'content': Regular text content lines (List[str])
                'lists': Bullet lists grouped (List[List[str]])
                'images': Image references with 'alt' and 'path'
                'speaker_notes': Speaker notes text (str, optional)
            base_path: Base directory path for resolving relative image paths.
                Used to resolve ./image.png style references.
            is_title_slide: If True, uses Title Slide layout (0) with centered title.
                If False, uses Title and Content layout (1) for title/body style.
                Defaults to False.

        Returns:
            None

        Side Effects:
            Adds a new slide to self.presentation with all parsed content.
            Modifies presentation state by adding shapes (text boxes, images).
            Adds speaker notes if present in slide_data.
            Applies background and font colors if configured.

        Examples:
            >>> converter = MarkdownToPowerPoint()
            >>> slide_data = {
            ...     "title": "My Slide",
            ...     "content": ["Some text"],
            ...     "lists": [["Item 1", "Item 2"]],
            ...     "images": [],
            ...     "speaker_notes": "Remember to mention key points",
            ... }
            >>> converter.add_slide_to_presentation(slide_data, is_title_slide=True)
            >>> len(converter.presentation.slides) == 1
            True
        """
        # Choose layout based on slide type
        if is_title_slide:
            # Layout 0: Title Slide (title centered in middle)
            slide_layout = self.presentation.slide_layouts[0]
        else:
            # Layout 1: Title and Content (title at top, body area below)
            slide_layout = self.presentation.slide_layouts[1]

        slide = self.presentation.slides.add_slide(slide_layout)

        # Apply background color based on slide type
        bg_color = self.title_bg_color if is_title_slide else self.background_color
        if bg_color:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color

        # Add background image if specified (add it first so other content appears on top)
        if self.background_image and os.path.exists(self.background_image):
            try:
                # Add background image to cover the entire slide
                slide.shapes.add_picture(
                    self.background_image,
                    Inches(0),  # Left position
                    Inches(0),  # Top position
                    width=Inches(10),  # Standard slide width
                    height=Inches(7.5),  # Standard slide height
                )
            except Exception as e:
                print(f"Warning: Could not add background image {self.background_image}: {e}")
        elif self.background_image:
            print(f"Warning: Background image not found: {self.background_image}")

        # Handle title based on slide type
        title_color = self.title_font_color if is_title_slide else self.font_color

        if is_title_slide:
            # For title slides, use the built-in title placeholder (centered)
            if slide_data["title"] and slide.shapes.title:
                self._apply_text_formatting(
                    slide.shapes.title.text_frame,
                    slide_data["title"],
                    font_size=32,
                    color=title_color,
                )
                # Set title to bold by default
                for paragraph in slide.shapes.title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if not run.font.name == "Courier New":  # Don't override code font
                            run.font.bold = True
            # Title slides typically don't have body content, but track position anyway
            top_position = Inches(4.0)
        else:
            # For content slides, use the title placeholder at the top
            if slide_data["title"] and slide.shapes.title:
                self._apply_text_formatting(
                    slide.shapes.title.text_frame,
                    slide_data["title"],
                    font_size=32,
                    color=title_color,
                )
                # Set title to bold by default
                for paragraph in slide.shapes.title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if not run.font.name == "Courier New":  # Don't override code font
                            run.font.bold = True
            # Start content below the title
            top_position = Inches(1.5)

        # Render body items in document order (new approach with "body" field)
        if "body" in slide_data and slide_data["body"]:
            for body_item in slide_data["body"]:
                if body_item["type"] == "content":
                    top_position = self._render_text_content(
                        slide,
                        body_item["text"],
                        body_item.get("content_type", "text"),
                        top_position,
                    )

                elif body_item["type"] == "list":
                    top_position = self._render_list_block(slide, body_item["items"], top_position)

                elif body_item["type"] == "table":
                    # Render table using the Phase 2 renderer which creates a native pptx table.
                    table = body_item["table"]
                    try:
                        rendered_height = self._render_table(slide, top_position, table)
                    except Exception:
                        # If rendering fails for any reason, fallback to a textual rendering height.
                        rendered_height = max(
                            (len(table.get("rows", [])) + (1 if table.get("has_header") else 0)) * 0.25,
                            0.5,
                        )
                    top_position = Inches(top_position.inches + rendered_height + 0.15)

        # Backward compatibility: render old-style content and lists if body is empty
        elif slide_data.get("content"):
            # Calculate height based on number of lines and content
            content_item_count = len(slide_data["content"])
            estimated_height = content_item_count * 0.35
            estimated_height = min(estimated_height, 4.0)
            estimated_height = max(estimated_height, 0.5)

            content_box = slide.shapes.add_textbox(Inches(0.5), top_position, Inches(9), Inches(estimated_height))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            # Process each content line with formatting based on type
            for i, content_line in enumerate(slide_data["content"]):
                if content_line.strip():
                    # Get content type (defaults to "text" for backward compatibility)
                    content_type = slide_data["content_types"][i] if i < len(slide_data["content_types"]) else "text"

                    if i == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()

                    # Set spacing based on type
                    if content_type.startswith("h"):
                        p.space_before = Pt(6)
                        p.space_after = Pt(3)
                    else:
                        p.space_before = Pt(3)
                        p.space_after = Pt(3)

                    # Apply formatting based on content type
                    segments = self._parse_markdown_formatting(content_line)
                    for segment in segments:
                        run = p.add_run()
                        run.text = segment["text"]
                        if segment["bold"]:
                            run.font.bold = True
                        if segment["italic"]:
                            run.font.italic = True
                        if segment["code"]:
                            run.font.name = "Courier New"

                        # Set font size based on content type
                        if content_type == "h3":
                            run.font.size = Pt(22)
                            run.font.bold = True
                        elif content_type == "h4":
                            run.font.size = Pt(20)
                            run.font.bold = True
                        elif content_type in ["h5", "h6"]:
                            run.font.size = Pt(18)
                            run.font.bold = True
                        else:
                            run.font.size = Pt(16)

                        if self.font_color:
                            run.font.color.rgb = self.font_color

            top_position = Inches(top_position.inches + estimated_height + 0.2)

            # Add lists with optimized spacing
            for list_items in slide_data.get("lists", []):
                list_height = max(len(list_items) * 0.35, 0.5)
                list_box = slide.shapes.add_textbox(Inches(1), top_position, Inches(8), Inches(list_height))
                list_frame = list_box.text_frame
                list_frame.clear()
                list_frame.word_wrap = True

                for i, item in enumerate(list_items):
                    if i == 0:
                        p = list_frame.paragraphs[0]
                    else:
                        p = list_frame.add_paragraph()

                    # Reduce spacing between list items
                    p.space_before = Pt(0)
                    p.space_after = Pt(3)

                    # Add bullet point and then formatted text
                    bullet_run = p.add_run()
                    bullet_run.text = "• "
                    bullet_run.font.size = Pt(16)
                    if self.font_color:
                        bullet_run.font.color.rgb = self.font_color

                    # Parse and apply markdown formatting to list item
                    segments = self._parse_markdown_formatting(item)
                    for segment in segments:
                        run = p.add_run()
                        run.text = segment["text"]
                        if segment["bold"]:
                            run.font.bold = True
                        if segment["italic"]:
                            run.font.italic = True
                        if segment["code"]:
                            run.font.name = "Courier New"
                        run.font.size = Pt(14)
                        if self.font_color:
                            run.font.color.rgb = self.font_color

                top_position = Inches(top_position.inches + list_height + 0.15)

        # Add code blocks
        for code_block in slide_data.get("code_blocks", []):
            top_position = self._render_code_block(slide, code_block, top_position)

        # Add images
        for image_info in slide_data["images"]:
            top_position = self._render_image(slide, image_info, base_path, top_position)

        # Add speaker notes if present
        if slide_data.get("speaker_notes"):
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data["speaker_notes"]

        # Remove unused placeholder shapes
        has_title = bool(slide_data.get("title"))
        # Include tables discovered in the new 'body' field as body content to ensure placeholders are cleaned appropriately
        has_body_content = bool(
            slide_data.get("content")
            or slide_data.get("lists")
            or slide_data.get("images")
            or slide_data.get("code_blocks")
            or any(item.get("type") == "table" for item in slide_data.get("body", []))
        )
        self._remove_unused_placeholders(slide, has_title, has_body_content)

    def convert(
        self,
        markdown_file: str,
        output_file: str,
        background_image: Optional[str] = None,
    ) -> None:
        """Convert a markdown file to PowerPoint presentation.

        Reads markdown content from input file, parses slides separated by '---',
        and generates a PowerPoint presentation with support for titles, lists,
        content, and optional background images.

        Args:
            markdown_file: Path to input markdown file to convert
            output_file: Path where output .pptx file will be saved
            background_image: Optional path to background image file for all slides.
                If provided and file exists, image will be added as background to
                each slide.

        Returns:
            None

        Raises:
            FileNotFoundError: If markdown_file cannot be read
            ValueError: If markdown content is empty or contains no valid slides
            IOError: If output file cannot be written

        Examples:
            >>> converter = MarkdownToPowerPoint()
            >>> converter.convert("slides.md", "output.pptx")
            Presentation saved to: output.pptx

            >>> converter = MarkdownToPowerPoint(background_image="bg.jpg")
            >>> converter.convert("slides.md", "output.pptx")
            Presentation saved to: output.pptx
        """
        # Set background image if provided
        if background_image:
            # Handle relative paths - resolve relative to current working directory, not markdown file
            if not os.path.isabs(background_image):
                background_image = os.path.abspath(background_image)
            self.background_image = background_image

        # Read markdown content
        with open(markdown_file, "r", encoding="utf-8") as f:
            markdown_content = f.read()

        # Parse slides
        slides_content = self.parse_markdown_slides(markdown_content)

        if not slides_content:
            raise ValueError("No slides found in markdown content")

        # Get base path for resolving relative image paths
        base_path = os.path.dirname(os.path.abspath(markdown_file))

        # Process each slide
        for index, slide_content in enumerate(slides_content):
            slide_data = self.parse_slide_content(slide_content)
            # First slide starting with single # is a title slide
            is_title_slide = index == 0 and slide_content.strip().startswith("# ")
            self.add_slide_to_presentation(slide_data, base_path, is_title_slide)

        # Save presentation
        self.presentation.save(output_file)
        print(f"Presentation saved to: {output_file}")


def create_presentation(cfg: Config) -> int:
    """Create a PowerPoint presentation from one or more markdown files.

    Supports three distinct usage modes for flexible file handling:

    Mode 1 - Input/output pair:
        Single input file with explicit output filename.
        Example: md2ppt create input.md output.pptx
        Creates: output.pptx in current directory

    Mode 2 - Single file, auto output:
        Single input file with auto-generated output name.
        Example: md2ppt create input.md
        Creates: input.pptx in same directory as input.md

    Mode 3 - Multiple files with output directory:
        Multiple input files processed to specified output directory.
        Example: md2ppt create a.md b.md --output ./presentations/
        Creates: presentations/a.pptx, presentations/b.pptx

    Args:
        cfg: Config dataclass instance containing configuration parameters:
            filenames (List[str]): List of input markdown file paths to process
            output_path (str): Output directory path for multi-file mode (Mode 3).
                Empty string for single-file modes.
            output_file (str): Explicit output filename for Mode 1.
                Empty string for Modes 2 and 3.
            background_path (str): Optional path to background image file.
                Will be added to all slides if file exists.
            verbose (bool): Enable verbose logging output
            debug (bool): Enable debug mode with detailed output

    Returns:
        int: Return code (0 on success)

    Raises:
        FileNotFoundError: If any input markdown file does not exist
        ValueError: If markdown content is empty or contains no valid slides
        IOError: If output file or directory cannot be created or written

    Examples:
        >>> from presenter.config import Config
        >>> from presenter.converter import create_presentation
        >>> cfg = Config(filenames=["slides.md"], output_file="output.pptx")
        >>> create_presentation(cfg)
        Presentation saved to: output.pptx
        0

        >>> cfg = Config(
        ...     filenames=["deck1.md", "deck2.md"], output_path="./presentations/", background_path="template.jpg"
        ... )
        >>> create_presentation(cfg)
        Presentation saved to: presentations/deck1.pptx
        Presentation saved to: presentations/deck2.pptx
        0
    """
    # Validate input files exist
    for filename in cfg.filenames:
        if not os.path.exists(filename):
            logger.error(f"Input file not found: {filename}")
            raise FileNotFoundError(f"Input file not found: {filename}")

    # Create output directory if specified and doesn't exist
    if cfg.output_path and not os.path.exists(cfg.output_path):
        os.makedirs(cfg.output_path, exist_ok=True)
        if cfg.verbose:
            logger.info(f"Created output directory: {cfg.output_path}")

    # Create output directory for explicit output file if needed (Mode 1)
    if cfg.output_file:
        output_dir = os.path.dirname(cfg.output_file)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)
            if cfg.verbose:
                logger.info(f"Created output directory: {output_dir}")

    # Prepare background image (validate once for all files)
    background_image = None
    if cfg.background_path:
        if os.path.exists(cfg.background_path):
            background_image = cfg.background_path
            if cfg.verbose:
                logger.info(f"Using background image: {background_image}")
        else:
            logger.warning(f"Background image not found: {cfg.background_path}")

    # Process each input file
    for filename in cfg.filenames:
        # Determine output filename based on mode
        if cfg.output_file:
            # Mode 1: Input/output pair - use explicit output filename
            output_file = cfg.output_file
            if cfg.verbose:
                logger.info(f"Converting {filename} -> {output_file}")
        elif cfg.output_path:
            # Mode 2: Multiple files with output directory
            base_name_only = os.path.basename(os.path.splitext(filename)[0])
            output_filename = base_name_only + ".pptx"
            output_file = os.path.join(cfg.output_path, output_filename)
            if cfg.verbose:
                logger.info(f"Converting {filename} -> {output_file}")
        else:
            # Mode 3: Single file, auto-generate output in same directory
            base_name = os.path.splitext(filename)[0]
            output_file = base_name + ".pptx"
            if cfg.verbose:
                logger.info(f"Converting {filename} -> {output_file}")

        # Create converter and process file
        converter = MarkdownToPowerPoint(
            background_image=background_image,
            background_color=cfg.background_color,
            font_color=cfg.font_color,
            title_bg_color=cfg.title_bg_color,
            title_font_color=cfg.title_font_color,
        )
        converter.convert(filename, output_file, background_image)

    return 0
